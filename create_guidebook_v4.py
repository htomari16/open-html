"""
新規顧客向けガイドブック v4（スタイル改善版）
- 背景: ライトグレー
- ヘッダー: 大番号 + セクションラベル + タイトル（参考スライドスタイル）
- カード: 白背景 + 細いカラーライン
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

COLOR_BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
COLOR_RED        = RGBColor(0xCC, 0x00, 0x00)
COLOR_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_DARK_GRAY  = RGBColor(0x44, 0x44, 0x44)
COLOR_MID_GRAY   = RGBColor(0x88, 0x88, 0x88)
COLOR_GREEN      = RGBColor(0x00, 0x88, 0x55)
COLOR_BLUE       = RGBColor(0x22, 0x55, 0xCC)
COLOR_ORANGE     = RGBColor(0xCC, 0x66, 0x00)
COLOR_PURPLE     = RGBColor(0x66, 0x22, 0xAA)
COLOR_WARM       = RGBColor(0xBB, 0x55, 0x00)
BG_SLIDE         = RGBColor(0xF0, 0xF2, 0xF4)  # コンテンツスライド背景（ライトグレー）

# 役割タグ定義 (label, color)
ROLE_SRA   = ("代行",   COLOR_RED)
ROLE_OWNER = ("貴社",   COLOR_BLUE)
ROLE_SR    = ("社労士", COLOR_GREEN)
ROLE_SUPP  = ("支援",   COLOR_WARM)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ============================================================
# ユーティリティ
# ============================================================

def add_rounded_rect(slide, left, top, width, height,
                     fill_color=None, line_color=None, line_width=None, radius=0.05):
    """角丸矩形を追加（5 = ROUNDED_RECTANGLE）"""
    shape = slide.shapes.add_shape(5, left, top, width, height)
    shape.adjustments[0] = radius
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def tb(slide, text, left, top, width, height,
       font_size=12, bold=False, color=COLOR_DARK_GRAY,
       align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Noto Sans JP"
    return txBox


def header(slide, section_no, title, note=""):
    # 大きな薄いグレーの番号（左端の装飾）
    tb(slide, section_no,
       Inches(0.05), Inches(-0.08), Inches(1.8), Inches(1.3),
       font_size=80, bold=True, color=RGBColor(0xDD, 0xDD, 0xDD))
    # セクションラベル（小さい赤文字）
    tb(slide, f"SECTION {section_no}",
       Inches(1.85), Inches(0.15), Inches(8.0), Inches(0.3),
       font_size=8.5, bold=True, color=COLOR_RED)
    # タイトル
    tb(slide, title,
       Inches(1.85), Inches(0.42), Inches(10.8), Inches(0.65),
       font_size=22, bold=True, color=COLOR_BLACK)
    # 赤いアンダーライン
    add_rect(slide, Inches(1.85), Inches(1.05), Inches(11.3), Cm(0.07), fill_color=COLOR_RED)
    if note:
        tb(slide, note, Inches(9.2), Inches(0.15), Inches(3.9), Inches(0.3),
           font_size=8, color=COLOR_MID_GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, text):
    tb(slide, text, Inches(0.2), Inches(7.0), Inches(13.0), Inches(0.4),
       font_size=8, color=COLOR_MID_GRAY)


def section_door(prs, no_str, title_str, desc_str, color=COLOR_RED):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, COLOR_BLACK)
    tb(slide, no_str,
       Inches(1.5), Inches(1.8), Inches(10.0), Inches(1.6),
       font_size=72, bold=True, color=color, align=PP_ALIGN.CENTER)
    tb(slide, title_str,
       Inches(1.5), Inches(3.5), Inches(10.0), Inches(0.9),
       font_size=32, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    tb(slide, desc_str,
       Inches(1.5), Inches(4.6), Inches(10.0), Inches(0.6),
       font_size=13, color=COLOR_MID_GRAY, align=PP_ALIGN.CENTER)


def role_chip(slide, left, top, role):
    """小さな役割タグ（label, color）を描画"""
    label, color = role
    w = Cm(1.2)
    h = Cm(0.36)
    add_rect(slide, left, top, w, h, fill_color=color)
    tb(slide, label, left, top - Cm(0.01), w, h + Cm(0.02),
       font_size=8.5, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)


def checkbox(slide, left, top, color=COLOR_RED):
    add_rect(slide, left, top + Cm(0.15), Cm(0.42), Cm(0.42),
             line_color=color, line_width=Pt(1.5))


def small_tag(slide, left, top, width, label, color):
    add_rect(slide, left, top, width, Cm(0.36), fill_color=color)
    tb(slide, label, left + Cm(0.08), top - Cm(0.01), width, Cm(0.38),
       font_size=8.5, bold=True, color=COLOR_WHITE)


# ============================================================
# プレゼンテーション生成
# ============================================================
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BL = prs.slide_layouts[6]  # blank


# ============================================================
# Slide 1: 表紙
# ============================================================
slide = prs.slides.add_slide(BL)
add_bg(slide, COLOR_WHITE)

add_rect(slide, 0, 0, Inches(4.2), SLIDE_H, fill_color=COLOR_BLACK)
add_rect(slide, Inches(4.2), 0, Inches(0.06), SLIDE_H, fill_color=COLOR_RED)

tb(slide, "SEVENRICH 会計事務所",
   Inches(0.4), Inches(0.5), Inches(3.6), Inches(0.6),
   font_size=12, bold=True, color=COLOR_WHITE)
tb(slide, "新規顧客向け\nガイドブック",
   Inches(0.4), Inches(1.4), Inches(3.6), Inches(1.9),
   font_size=27, bold=True, color=COLOR_WHITE)
tb(slide, "会社設立後に知っておくべき\n税務・経理の基本",
   Inches(0.4), Inches(3.5), Inches(3.6), Inches(1.0),
   font_size=11, color=COLOR_MID_GRAY)
tb(slide, "設立おめでとうございます。\nともに歩んでまいりましょう。",
   Inches(0.4), Inches(4.7), Inches(3.6), Inches(0.9),
   font_size=10, italic=True, color=RGBColor(0xAA, 0xAA, 0xAA))
tb(slide, "2026年版",
   Inches(0.4), Inches(6.7), Inches(3.6), Inches(0.4),
   font_size=9, color=COLOR_MID_GRAY)

tb(slide, "CONTENTS",
   Inches(4.6), Inches(0.55), Inches(8.0), Inches(0.45),
   font_size=11, bold=True, color=COLOR_RED)

contents = [
    ("01", "設立後の対応チェックリスト",    "設立直後〜3ヶ月以内に行う各種手続き（役割分担付き）"),
    ("02", "月次の社内業務チェックリスト",   "まず3つのルールから始める経理・労務業務"),
    ("03", "年間の税務スケジュール",         "申告・納付カレンダー（決算月基準の相対表示）"),
    ("04", "消費税の基本",                   "免税・課税の判定とBtoB/BtoC別の目安"),
    ("05", "役員報酬の基本",                 "3ヶ月ルールと金額の決め方ガイド"),
    ("06", "会社と個人のお金の分離",         "一番大事な基本ルール―ここだけは絶対に守る"),
    ("07", "法人口座・クレジットカード",         "おすすめ銀行・カードの選び方と4口座の使い分け"),
    ("08", "SEVENRICHとのやり取り＆SOSリスト", "連絡方法・提出期限・動く前に相談すべきこと"),
    ("09", "よくある質問（FAQ）",            "新規顧客からよく頂く質問と回答"),
]

for i, (no, title, desc) in enumerate(contents):
    y = Inches(1.1) + i * Inches(0.63)
    add_rect(slide, Inches(4.6), y, Inches(8.4), Inches(0.60), fill_color=COLOR_LIGHT_GRAY)
    tb(slide, no,    Inches(4.75), y + Inches(0.04), Inches(0.55), Inches(0.33),
       font_size=14, bold=True, color=COLOR_RED)
    tb(slide, title, Inches(5.45), y + Inches(0.03), Inches(7.3), Inches(0.30),
       font_size=12, bold=True, color=COLOR_BLACK)
    tb(slide, desc,  Inches(5.45), y + Inches(0.33), Inches(7.3), Inches(0.24),
       font_size=8, color=COLOR_DARK_GRAY)


# ============================================================
# Slide 2: ウェルカムスライド（NEW）
# ============================================================
slide = prs.slides.add_slide(BL)
add_bg(slide, COLOR_BLACK)

# 大きな設立おめでとう
tb(slide, "設立おめでとうございます！",
   Inches(1.0), Inches(0.6), Inches(11.3), Inches(1.1),
   font_size=36, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
tb(slide, "これからの経営を、しっかりとサポートさせていただきます。手続きや数字のことは、どうぞ私どもにお任せください。",
   Inches(1.5), Inches(1.7), Inches(10.3), Inches(0.6),
   font_size=13, color=COLOR_MID_GRAY, align=PP_ALIGN.CENTER)

# 3つの約束
promises = [
    ("01", "手続きはお任せください",
     "設立直後に必要な税務署・自治体への届出は、当事務所が代行いたします。\n社長様はビジネスに専念いただけます。"),
    ("02", "毎月の数字を一緒に見ます",
     "月次の試算表を毎月確認し、資金繰りや節税のアドバイスをお伝えします。\n「数字がわからない」というお悩みは、一緒に解消してまいります。"),
    ("03", "ご実施前にご相談ください",
     "役員報酬の変更・大きなお買い物・採用など、ご実施前にひとことご相談いただけますと、\n最大限の税務メリットをご提案することができます。"),
]

for i, (no, title, desc) in enumerate(promises):
    left = Inches(0.4) + i * Inches(4.3)
    top  = Inches(2.5)
    w    = Inches(4.0)
    h    = Inches(4.3)
    add_rect(slide, left, top, w, h, fill_color=RGBColor(0x28, 0x28, 0x28))
    add_rect(slide, left, top, w, Inches(0.12), fill_color=COLOR_RED)
    tb(slide, no,    left, top + Inches(0.2), w, Inches(0.6),
       font_size=28, bold=True, color=COLOR_RED, align=PP_ALIGN.CENTER)
    tb(slide, title, left + Inches(0.2), top + Inches(0.9), w - Inches(0.4), Inches(0.55),
       font_size=15, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    tb(slide, desc,  left + Inches(0.2), top + Inches(1.55), w - Inches(0.4), Inches(2.4),
       font_size=10, color=COLOR_MID_GRAY, align=PP_ALIGN.LEFT)

tb(slide, "このガイドブックをご活用いただきながら、最初の1年を一緒に歩んでまいりましょう。",
   Inches(1.5), Inches(6.9), Inches(10.3), Inches(0.45),
   font_size=11, italic=True, color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)


# ============================================================
# Slides 3-4: 01 設立後チェックリスト（役割タグ付き）
# ============================================================
section_door(prs, "01", "設立後の対応チェックリスト",
             "「誰が・いつまでに」を明確にした手続き一覧です")

slide = prs.slides.add_slide(BL)
add_bg(slide, BG_SLIDE)
header(slide, "01", "設立後の対応チェックリスト", "設立直後〜3ヶ月以内")

# 役割凡例
legend_top = Inches(0.95)
legend_x = Inches(9.5)
tb(slide, "役割：", legend_x, legend_top, Inches(0.5), Inches(0.28),
   font_size=8, color=COLOR_MID_GRAY)
for i, (label, color) in enumerate([ROLE_SRA, ROLE_OWNER, ROLE_SR, ROLE_SUPP]):
    add_rect(slide, legend_x + Inches(0.45) + i * Inches(1.05), legend_top + Cm(0.03),
             Cm(1.2), Cm(0.36), fill_color=color)
    tb(slide, label,
       legend_x + Inches(0.45) + i * Inches(1.05), legend_top - Cm(0.01),
       Cm(1.2), Cm(0.38),
       font_size=8.5, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

col_w   = Inches(3.9)
col_h   = Inches(5.85)
col_top = Inches(1.15)
col_lefts = [Inches(0.2), Inches(4.3), Inches(8.4)]

# (item_title, item_desc, role)
check_sections = [
    {
        "title": "A. 税務署・自治体への届出",
        "color": COLOR_RED,
        "items": [
            ("法人設立届出書",            "税務署・都道府県・市区町村 各1通",      ROLE_SRA),
            ("青色申告の承認申請書",      "設立後3ヶ月以内または最初の期末まで",    ROLE_SRA),
            ("給与支払事務所等の開設届",  "開設から1ヶ月以内",                      ROLE_SRA),
            ("源泉所得税 納期の特例申請", "常時10人未満の場合に申請可",             ROLE_SRA),
            ("消費税 課税事業者選択届",   "インボイス登録と合わせて要相談",          ROLE_SRA),
            ("棚卸資産の評価方法の届出",  "最初の期の確定申告期限まで",              ROLE_SRA),
            ("減価償却資産の償却方法届",  "同上",                                    ROLE_SRA),
        ]
    },
    {
        "title": "B. 金融機関・会計準備",
        "color": COLOR_BLUE,
        "items": [
            ("法人口座の開設",            "設立登記後すぐにご開設ください（審査にお時間がかかります）",  ROLE_OWNER),
            ("法人クレジットカードの申請","経費管理の効率化のため、お早めにお申請ください",            ROLE_OWNER),
            ("会計ソフトの導入・初期設定","MFクラウド会計 / freee 等（初期設定は当事務所がサポートいたします）",              ROLE_SUPP),
            ("電子帳簿保存法の対応確認",  "電子取引データの保存方法を整備",          ROLE_SUPP),
            ("インボイス登録申請",         "登録番号の取得（適格請求書発行事業者）",  ROLE_SRA),
            ("法人印鑑証明書の取得",       "各種手続きに複数枚必要",                 ROLE_OWNER),
            ("許認可の取得",               "業種による（行政書士へ依頼）",            ROLE_OWNER),
        ]
    },
    {
        "title": "C. 社会保険・労務",
        "color": COLOR_GREEN,
        "items": [
            ("健保・厚生年金 新規適用届",  "設立から5日以内",                        ROLE_SR),
            ("雇用保険 適用事業所設置届",  "設立翌日から10日以内",                   ROLE_SR),
            ("労災保険 保険関係成立届",    "保険関係成立翌日から10日以内",           ROLE_SR),
            ("役員報酬の決定",             "設立後3ヶ月以内（定期同額給与の要件）",  ROLE_OWNER),
            ("就業規則の作成",             "常時10人以上は届出義務あり",              ROLE_SR),
            ("雇用契約書・労働条件通知書", "雇用開始前に交付が必要",                 ROLE_SR),
            ("給与規程・旅費規程の整備",   "経費精算ルールの明文化",                  ROLE_SUPP),
        ]
    },
]

ITEM_H = Inches(0.75)

for left, sec in zip(col_lefts, check_sections):
    add_rounded_rect(slide, left, col_top, col_w, col_h, fill_color=COLOR_WHITE)
    add_rect(slide, left, col_top, col_w, Inches(0.07), fill_color=sec["color"])
    tb(slide, sec["title"],
       left + Cm(0.3), col_top + Inches(0.1), col_w, Inches(0.4),
       font_size=10.5, bold=True, color=sec["color"])

    for i, (t, d, role) in enumerate(sec["items"]):
        iy = col_top + Inches(0.5) + i * ITEM_H
        # チェックボックス
        checkbox(slide, left + Cm(0.25), iy, sec["color"])
        # 役割タグ
        rl, rc = role
        add_rect(slide, left + Cm(1.0), iy + Cm(0.08), Cm(1.2), Cm(0.36), fill_color=rc)
        tb(slide, rl, left + Cm(1.0), iy + Cm(0.06), Cm(1.2), Cm(0.38),
           font_size=8.5, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
        # タイトル
        tb(slide, t,
           left + Cm(2.2), iy + Cm(0.03), col_w - Cm(2.4), Inches(0.3),
           font_size=9.5, bold=True, color=COLOR_BLACK)
        # 説明
        tb(slide, d,
           left + Cm(1.0), iy + Cm(0.45), col_w - Cm(1.2), Inches(0.32),
           font_size=7.5, color=COLOR_MID_GRAY)

footer(slide, "※ 届出書の提出期限は設立日基準のものが多く、遅延すると不利益が生じます。提出は当事務所が代行いたします（C列は社労士へご依頼いたします）。")


# ============================================================
# Slides 5-6: 02 月次業務チェックリスト（「まず3つ！」追加）
# ============================================================
section_door(prs, "02", "月次の社内業務チェックリスト",
             "まず3つのルールから始めてください。あとはSRAがサポートします。")

slide = prs.slides.add_slide(BL)
add_bg(slide, BG_SLIDE)
header(slide, "02", "月次の社内業務チェックリスト")

# 「まず3つだけ！」ハイライトボックス
hb_top = Inches(1.05)
add_rect(slide, Inches(0.2), hb_top, Inches(13.0), Inches(1.05), fill_color=RGBColor(0xFF, 0xF0, 0xF0))
add_rect(slide, Inches(0.2), hb_top, Inches(0.1), Inches(1.05), fill_color=COLOR_RED)
tb(slide, "まず、これだけ！3つのルール（仕訳は当事務所がチェックいたします。社長様はこれだけご意識ください）",
   Inches(0.5), hb_top + Inches(0.04), Inches(12.5), Inches(0.38),
   font_size=11, bold=True, color=COLOR_RED)
rules3 = [
    "① 領収書は捨てずに箱 or MFクラウドで写真撮影してください",
    "② 会社の経費は個人のお財布からお支払いにならないようにしてください",
    "③ 月末に取引先への請求書を必ずご発行ください",
]
tb(slide, "    ".join(rules3),
   Inches(0.5), hb_top + Inches(0.48), Inches(12.5), Inches(0.45),
   font_size=10, bold=False, color=COLOR_DARK_GRAY)

# 4カラム（日次/週次/月初/月末）
col_w2   = Inches(2.9)
col_h2   = Inches(4.45)
col_top2 = Inches(2.25)
col_l2   = [Inches(0.2), Inches(3.4), Inches(6.4), Inches(9.5)]

monthly = [
    {
        "title": "日次",
        "color": COLOR_PURPLE,
        "items": [
            ("領収書・請求書の収集",    "紙・電子ともに当日中に保管"),
            ("現金出納帳の記帳",        "現金残高を実際に確認"),
            ("銀行入出金の確認",        "ネットバンキングで毎日チェック"),
            ("売上・入金の記録",        "入金漏れを早期発見するため"),
        ]
    },
    {
        "title": "週次",
        "color": COLOR_BLUE,
        "items": [
            ("仕訳入力（会計ソフト）",  "週1回まとめて入力が目安"),
            ("未払請求書の確認",        "支払期日の管理"),
            ("売掛金・買掛金の確認",    "回収漏れ・支払遅延の防止"),
            ("経費精算の処理",          "社員からの申請をまとめて処理"),
        ]
    },
    {
        "title": "月初（1〜10日）",
        "color": COLOR_RED,
        "items": [
            ("前月の仕訳締め・確認",    "会計ソフトで月次締め処理"),
            ("給与計算",                "締め日・支払日に合わせて実施"),
            ("源泉所得税の納付",        "翌月10日まで（特例：7月・1月）"),
            ("社会保険料の納付",        "翌月末まで（口座振替推奨）"),
            ("試算表の確認",            "損益・貸借をSRAと共に確認"),
            ("資金繰り表の更新",        "向こう3ヶ月の資金計画を見直し"),
        ]
    },
    {
        "title": "月末",
        "color": COLOR_ORANGE,
        "items": [
            ("売上・費用の計上確認",    "発生主義で漏れなく計上"),
            ("請求書の発行・送付",      "取引先への請求を月末締めで"),
            ("預金残高の突合",          "帳簿残高と通帳残高の一致確認"),
            ("棚卸",                    "在庫がある場合は月末に実施"),
            ("翌月の資金繰り確認",      "支払予定・入金予定の洗い出し"),
            ("月次レポートの共有",      "MFクラウドでSRAへ資料連携"),
        ]
    },
]

ITEM_H2 = Inches(0.62)

for left, sec in zip(col_l2, monthly):
    add_rounded_rect(slide, left, col_top2, col_w2, col_h2, fill_color=COLOR_WHITE)
    add_rect(slide, left, col_top2, col_w2, Inches(0.07), fill_color=sec["color"])
    tb(slide, sec["title"],
       left + Cm(0.2), col_top2 + Inches(0.1), col_w2, Inches(0.4),
       font_size=12, bold=True, color=sec["color"], align=PP_ALIGN.CENTER)
    for i, (t, d) in enumerate(sec["items"]):
        iy = col_top2 + Inches(0.5) + i * ITEM_H2
        checkbox(slide, left + Cm(0.25), iy, sec["color"])
        tb(slide, t,
           left + Cm(0.95), iy, col_w2 - Cm(1.1), Inches(0.28),
           font_size=9, bold=True, color=COLOR_BLACK)
        tb(slide, d,
           left + Cm(0.95), iy + Inches(0.27), col_w2 - Cm(1.1), Inches(0.3),
           font_size=7.5, color=COLOR_MID_GRAY)

footer(slide, "※ 源泉所得税の納期の特例（常時10人未満）をご選択の場合、1〜6月分は7月10日、7〜12月分は翌年1月20日に一括納付いただきます。")


# ============================================================
# Slides 7-8: 03 年間税務スケジュール（相対表示 M+N）
# ============================================================
section_door(prs, "03", "年間の税務スケジュール",
             "M = 貴社の決算月。M+1が翌月、M+2が翌々月…と置き換えて読んでください。")

slide = prs.slides.add_slide(BL)
add_bg(slide, BG_SLIDE)
header(slide, "03", "年間の税務スケジュール", "M = 決算月  ★ = まとまった出費が発生します")

# ★説明凡例
add_rect(slide, Inches(10.4), Inches(0.15), Inches(2.7), Inches(0.55), fill_color=RGBColor(0xFF, 0xEE, 0xEE))
tb(slide, "★ = 事前に資金を確保してください",
   Inches(10.5), Inches(0.18), Inches(2.5), Inches(0.5),
   font_size=8.5, bold=True, color=COLOR_RED)

# 相対月イベント
rel_months = [
    ("M+1",          False, [("源泉所得税", COLOR_BLUE,   "前月（期末月）分 納付（10日まで）")]),
    ("M+2\n申告期限", True,  [("法人税・地方税", COLOR_RED,  "★ 申告・納付（期末後2ヶ月以内）"),
                               ("消費税",     COLOR_ORANGE,"★ 申告・納付（同上）"),
                               ("源泉所得税", COLOR_BLUE,  "前月分 納付（10日）")]),
    ("M+3",          False, [("源泉所得税", COLOR_BLUE,   "前月分 納付（10日まで）")]),
    ("M+4",          False, [("源泉所得税\n（特例：7月）", COLOR_BLUE, "1〜6月分 一括納付（20日）"),
                              ("源泉所得税\n（特例以外）", COLOR_BLUE, "前月分 納付（10日）")]),
    ("M+5",          False, [("源泉所得税", COLOR_BLUE,   "前月分 納付（10日まで）")]),
    ("M+6\n中間申告", True,  [("法人税等",  COLOR_RED,    "★ 中間申告・納付（前期税額の1/2）"),
                               ("消費税",     COLOR_ORANGE,"中間申告（年1回中間の場合）"),
                               ("源泉所得税", COLOR_BLUE,  "前月分 納付（10日）")]),
    ("M+7",          False, [("源泉所得税", COLOR_BLUE,   "前月分 納付（10日まで）")]),
    ("M+8",          False, [("源泉所得税", COLOR_BLUE,   "前月分 納付"),
                              ("消費税",     COLOR_ORANGE, "中間申告（年3回の場合・2回目）")]),
    ("M+9",          False, [("源泉所得税", COLOR_BLUE,   "前月分 納付"),
                              ("法人税等",   COLOR_RED,    "中間申告（年3回の場合・3回目）")]),
    ("M+10",         False, [("源泉所得税\n（特例：1月）", COLOR_BLUE, "7〜12月分 一括納付（20日）"),
                              ("源泉所得税\n（特例以外）", COLOR_BLUE, "前月分 納付（10日）")]),
    ("M+11",         False, [("年末調整",   COLOR_GREEN,  "12月給与で年末調整を実施（毎年12月固定）"),
                              ("源泉所得税", COLOR_BLUE,   "前月分 納付")]),
    ("M+12\n（決算月）", False, [("決算",      COLOR_RED,    "期末 → 決算資料の収集・整理開始"),
                                  ("法定調書",  COLOR_PURPLE, "給与支払報告書・合計表（1月31日固定）"),
                                  ("源泉所得税",COLOR_BLUE,   "前月分 納付")]),
]

col_count = 4
cell_w = Inches(3.1)
cell_h = Inches(1.72)
s_left = Inches(0.2)
s_top  = Inches(1.05)

for idx, (month, is_warn, tasks) in enumerate(rel_months):
    row = idx // col_count
    col = idx % col_count
    left = s_left + col * (cell_w + Inches(0.08))
    top  = s_top  + row * (cell_h + Inches(0.05))

    bg = RGBColor(0xFF, 0xF2, 0xF2) if is_warn else COLOR_WHITE
    add_rounded_rect(slide, left, top, cell_w, cell_h, fill_color=bg, radius=0.03)
    lbl_color = COLOR_RED if is_warn else COLOR_BLACK
    add_rect(slide, left, top, Inches(0.72), cell_h, fill_color=lbl_color)
    tb(slide, month,
       left, top, Inches(0.72), cell_h,
       font_size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

    for t_idx, (tax_type, t_color, detail) in enumerate(tasks):
        ty = top + Inches(0.06) + t_idx * Inches(0.52)
        small_tag(slide, left + Inches(0.78), ty + Cm(0.05), Inches(1.0), tax_type.split("\n")[0], t_color)
        tb(slide, detail,
           left + Inches(0.8), ty + Cm(0.36), Inches(2.2), Inches(0.42),
           font_size=7.5, color=COLOR_DARK_GRAY)

footer(slide, "※ 年末調整（毎年12月）・法定調書（毎年1月31日）・償却資産税申告（毎年1月31日）は暦年固定のイベントです。M+Nとは別にご管理ください。")


# ============================================================
# Slides 9-10: 04 消費税の基本（BtoB/BtoC 目安追加）
# ============================================================
section_door(prs, "04", "消費税の基本",
             "免税事業者のルールと、インボイス制度対応の考え方")

slide = prs.slides.add_slide(BL)
add_bg(slide, BG_SLIDE)
header(slide, "04", "消費税の基本")

# 左：判定フロー
L, LW, T = Inches(0.2), Inches(5.9), Inches(1.1)

tb(slide, "免税事業者の判定フロー",
   L, T, LW, Inches(0.4), font_size=12, bold=True, color=COLOR_BLACK)

flow_steps = [
    ("STEP 1", "基準期間（2期前）の課税売上高は？",
     "1,000万円以下 → 次のSTEPへ　/　超 → 課税事業者（申告義務あり）"),
    ("STEP 2", "特定期間（前期上半期）の売上・給与は？",
     "各1,000万円以下 → 次のSTEPへ　/　超 → 課税事業者"),
    ("STEP 3", "資本金は1,000万円未満か？",
     "1,000万円以上 → 設立1・2期目から課税事業者になる"),
    ("STEP 4", "インボイス（適格請求書）の登録をするか？",
     "登録する → 登録日から課税事業者（免税期間でも課税）"),
]

for i, (step, q, ans) in enumerate(flow_steps):
    fy = T + Inches(0.5) + i * Inches(1.3)
    add_rounded_rect(slide, L, fy, LW, Inches(1.18), fill_color=COLOR_WHITE, radius=0.03)
    add_rect(slide, L, fy, Inches(0.82), Inches(1.18), fill_color=COLOR_ORANGE)
    tb(slide, step,
       L, fy, Inches(0.82), Inches(1.18),
       font_size=8, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    tb(slide, q,
       L + Inches(0.9), fy + Inches(0.1), LW - Inches(1.0), Inches(0.4),
       font_size=10, bold=True, color=COLOR_BLACK)
    tb(slide, ans,
       L + Inches(0.9), fy + Inches(0.55), LW - Inches(1.0), Inches(0.55),
       font_size=8.5, color=COLOR_DARK_GRAY)

# 右：BtoB/BtoC目安 + 届出書
R, RW = Inches(6.4), Inches(6.7)

tb(slide, "登録すべき？ ケース別の目安",
   R, T, RW, Inches(0.4), font_size=12, bold=True, color=COLOR_BLACK)

cases = [
    {
        "tag": "BtoB 中心",
        "color": COLOR_RED,
        "body": "取引先が法人・個人事業主メインの場合\n→ インボイス登録（課税事業者）がほぼ必須\n  未登録だと相手が仕入税額控除できず、取引を敬遠される可能性あり",
        "result": "登録を強く推奨",
    },
    {
        "tag": "BtoC 中心",
        "color": COLOR_GREEN,
        "body": "エンドユーザー（個人消費者）がメインの場合\n→ 売上が1,000万円以下なら免税のままでも可\n  ただし将来的に売上が伸びる場合はタイミングを検討",
        "result": "免税継続も選択肢",
    },
    {
        "tag": "届出書と期限",
        "color": COLOR_ORANGE,
        "body": "課税事業者選択届：適用課税期間の開始前日\n簡易課税制度選択届：同上（売上5,000万円以下が対象）\n消費税の中間申告：前期税額48万円超で義務",
        "result": "期限厳守",
    },
]

for i, case in enumerate(cases):
    cy = T + Inches(0.5) + i * Inches(1.72)
    add_rounded_rect(slide, R, cy, RW, Inches(1.55), fill_color=COLOR_WHITE, radius=0.03)
    add_rect(slide, R, cy, Inches(0.2), Inches(1.55), fill_color=case["color"])
    add_rect(slide, R + Inches(0.3), cy + Inches(0.1),
             Inches(1.4), Cm(0.38), fill_color=case["color"])
    tb(slide, case["tag"],
       R + Inches(0.3), cy + Inches(0.08), Inches(1.4), Cm(0.40),
       font_size=9, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    tb(slide, case["body"],
       R + Inches(0.3), cy + Inches(0.44), RW - Inches(0.5), Inches(0.8),
       font_size=8.5, color=COLOR_DARK_GRAY)
    add_rect(slide, R + RW - Inches(1.4), cy + Inches(1.12),
             Inches(1.2), Cm(0.3), fill_color=case["color"])
    tb(slide, case["result"],
       R + RW - Inches(1.4), cy + Inches(1.1), Inches(1.2), Cm(0.32),
       font_size=7.5, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

footer(slide, "ご不明な点はお気軽にご相談ください。登録の要否は、売上状況・取引先構成を踏まえ、一緒に検討いたします。")


# ============================================================
# Slides 11-12: 05 役員報酬の基本（金額設定ガイド追加）
# ============================================================
section_door(prs, "05", "役員報酬の基本",
             "3ヶ月ルールを守りつつ、最適な金額を一緒に設計します")

slide = prs.slides.add_slide(BL)
add_bg(slide, BG_SLIDE)
header(slide, "05", "役員報酬の基本")

# 左：3種類の説明
L, LW, T = Inches(0.2), Inches(5.9), Inches(1.1)

tb(slide, "役員報酬の3種類",
   L, T, LW, Inches(0.4), font_size=12, bold=True, color=COLOR_BLACK)

types = [
    ("定期同額給与",     COLOR_RED,
     "毎月同額を支払う。最も一般的な方法。\n変更できるのは「事業年度開始から3ヶ月以内」のみ。期間外の変更分は損金不算入。"),
    ("事前確定届出給与", COLOR_BLUE,
     "賞与相当。支給額・日付を事前に税務署へ届け出ることで損金算入可。\n届出期限：定時株主総会から1ヶ月以内 または 事業年度開始から4ヶ月以内。"),
    ("業績連動給与",     COLOR_PURPLE,
     "利益連動型。同族会社は原則不可。上場企業向けの特殊制度。"),
]

for i, (t_title, t_color, t_desc) in enumerate(types):
    ty = T + Inches(0.5) + i * Inches(1.65)
    h  = Inches(1.5)
    add_rounded_rect(slide, L, ty, LW, h, fill_color=COLOR_WHITE, radius=0.03)
    add_rect(slide, L, ty, Inches(0.2), h, fill_color=t_color)
    tb(slide, t_title,
       L + Inches(0.3), ty + Inches(0.1), LW - Inches(0.4), Inches(0.38),
       font_size=11, bold=True, color=COLOR_BLACK)
    tb(slide, t_desc,
       L + Inches(0.3), ty + Inches(0.52), LW - Inches(0.4), Inches(0.88),
       font_size=8.5, color=COLOR_DARK_GRAY)

# 右：金額の決め方 + 節税ポイント
R, RW = Inches(6.4), Inches(6.7)

tb(slide, "金額はどう決める？ 設定ガイド",
   R, T, RW, Inches(0.4), font_size=12, bold=True, color=COLOR_BLACK)

guides = [
    {
        "step": "STEP 1",
        "title": "生活に必要な最低額から逆算する",
        "desc": "手取りで月〇〇万円必要 → 所得税・住民税・社会保険を逆算し、\n総支給額（役員報酬）を算出します。生活費がベースラインです。",
    },
    {
        "step": "STEP 2",
        "title": "社会保険料の「壁」を意識する",
        "desc": "標準報酬月額は等級ごとに上限あり（最高139万円）。\n報酬が上がるほど社会保険料も増えるため、法人・個人の合算でシミュレーションします。",
    },
    {
        "step": "STEP 3",
        "title": "法人の利益計画と照らし合わせる",
        "desc": "役員報酬を増やす → 法人税↓・個人所得税↑。\n最適なバランスはその年の利益予測によって変わるため、期初に一緒に試算します。",
    },
    {
        "step": "提案",
        "title": "初回キックオフ面談でシミュレーションいたします",
        "desc": "設立後すぐに「利益計画×役員報酬シミュレーション」を作成いたします。\n「金額の決め方がわからない」というお悩みも、一緒に解決いたします。",
    },
]

for i, g in enumerate(guides):
    gy = T + Inches(0.5) + i * Inches(1.3)
    is_proposal = g["step"] == "提案"
    bg = RGBColor(0xFF, 0xF0, 0xF0) if is_proposal else COLOR_LIGHT_GRAY
    sc = COLOR_RED if is_proposal else COLOR_DARK_GRAY
    add_rounded_rect(slide, R, gy, RW, Inches(1.15), fill_color=bg, radius=0.03)
    add_rect(slide, R, gy, Inches(0.9), Inches(1.15), fill_color=sc)
    tb(slide, g["step"],
       R, gy, Inches(0.9), Inches(1.15),
       font_size=9, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    tb(slide, g["title"],
       R + Inches(1.0), gy + Inches(0.1), RW - Inches(1.1), Inches(0.38),
       font_size=10, bold=True, color=COLOR_BLACK)
    tb(slide, g["desc"],
       R + Inches(1.0), gy + Inches(0.5), RW - Inches(1.1), Inches(0.58),
       font_size=8.5, color=COLOR_DARK_GRAY)

footer(slide, "※ 役員報酬の変更は期初3ヶ月以内が鉄則です。変更をご検討の際は、必ず事前にご相談ください。")


# ============================================================
# Slides 13-14: 06 会社と個人のお金の分離（NEW）
# ============================================================
section_door(prs, "06", "会社と個人のお金の分離",
             "これが一番大事な基本ルール。最初に必ず押さえてください。",
             color=COLOR_ORANGE)

slide = prs.slides.add_slide(BL)
add_bg(slide, BG_SLIDE)
header(slide, "06", "会社と個人のお金の分離")

# NG vs OK の2カラム
MID = Inches(6.7)

# NG列
add_rect(slide, Inches(0.2), Inches(1.05), MID - Inches(0.4), Inches(6.0),
         fill_color=RGBColor(0xFF, 0xF0, 0xF0))
add_rect(slide, Inches(0.2), Inches(1.05), MID - Inches(0.4), Inches(0.55),
         fill_color=RGBColor(0xBB, 0x00, 0x00))
tb(slide, "NG  やってはいけない行動",
   Inches(0.4), Inches(1.05), MID - Inches(0.6), Inches(0.55),
   font_size=13, bold=True, color=COLOR_WHITE)

# OK列
add_rect(slide, MID + Inches(0.1), Inches(1.05), MID - Inches(0.5), Inches(6.0),
         fill_color=RGBColor(0xF0, 0xFF, 0xF5))
add_rect(slide, MID + Inches(0.1), Inches(1.05), MID - Inches(0.5), Inches(0.55),
         fill_color=COLOR_GREEN)
tb(slide, "OK  正しいやり方",
   MID + Inches(0.3), Inches(1.05), MID - Inches(0.6), Inches(0.55),
   font_size=13, bold=True, color=COLOR_WHITE)

ng_items = [
    ("個人カードで会社の経費を払う",
     "何が会社の経費か追えなくなります。税務調査で否認されるリスクが高まります。"),
    ("会社口座から生活費を引き出す",
     "「役員貸付金」となり法人税上の問題が生じます。最悪の場合、認定利息課税も発生します。"),
    ("個人と会社の口座を同じ通帳で管理",
     "決算処理が複雑になり、正確な損益計算ができなくなります。"),
    ("個人名義の固定費を会社費用にする",
     "按分や根拠がない場合、全額否認されることがあります。事前に必ずご相談ください。"),
    ("レシート・領収書をすぐ捨てる",
     "経費の証拠がなくなります。最低7年間の保存義務がございます。"),
]

ok_items = [
    ("法人口座・法人カードを設立直後に開設",
     "経費はすべて法人カードでお支払いください。通帳は会社専用を維持してください。"),
    ("生活費は毎月の役員報酬から賄う",
     "報酬以外で個人口座へのご入金はお控えください。不足する場合は役員報酬額の見直しをご相談ください。"),
    ("立替払いをした場合は「社員立替精算書」で処理",
     "やむを得ず個人払いをされた場合は記録を残し、月次でご精算ください。"),
    ("個人・法人の費用を按分するときは根拠を記録",
     "自宅兼事務所・車両等は使用割合を計算し、書面でお残しください。"),
    ("毎月MFクラウドで帳簿を最新状態に保つ",
     "リアルタイムの帳簿は税務調査対策にも、資金繰り管理にも大変有効です。"),
]

NG_C = RGBColor(0xBB, 0x00, 0x00)
OK_C = COLOR_GREEN

for i, (t, d) in enumerate(ng_items):
    iy = Inches(1.72) + i * Inches(0.99)
    add_rect(slide, Inches(0.35), iy + Cm(0.1), Cm(0.55), Cm(0.55),
             fill_color=NG_C)
    tb(slide, "×", Inches(0.35), iy + Cm(0.06), Cm(0.55), Cm(0.6),
       font_size=13, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    tb(slide, t,
       Inches(0.9), iy, MID - Inches(1.0), Inches(0.35),
       font_size=9.5, bold=True, color=NG_C)
    tb(slide, d,
       Inches(0.9), iy + Inches(0.34), MID - Inches(1.0), Inches(0.45),
       font_size=8, color=COLOR_DARK_GRAY)

for i, (t, d) in enumerate(ok_items):
    iy = Inches(1.72) + i * Inches(0.99)
    add_rect(slide, MID + Inches(0.25), iy + Cm(0.1), Cm(0.55), Cm(0.55),
             fill_color=OK_C)
    tb(slide, "○", MID + Inches(0.25), iy + Cm(0.06), Cm(0.55), Cm(0.6),
       font_size=13, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    tb(slide, t,
       MID + Inches(0.8), iy, MID - Inches(1.1), Inches(0.35),
       font_size=9.5, bold=True, color=OK_C)
    tb(slide, d,
       MID + Inches(0.8), iy + Inches(0.34), MID - Inches(1.1), Inches(0.45),
       font_size=8, color=COLOR_DARK_GRAY)

footer(slide, "公私混合が続くと決算処理が複雑になり、追加費用が発生します。最初の習慣づけが最も重要です。")


# ============================================================
# Slides 15-16: 07 法人口座・クレジットカード
# ============================================================
section_door(prs, "07", "法人口座・クレジットカード",
             "設立後すぐに開設！口座4種類の使い分けとおすすめ銀行・カード",
             color=COLOR_BLUE)

slide = prs.slides.add_slide(BL)
add_bg(slide, BG_SLIDE)
header(slide, "07", "法人口座・クレジットカード")

# ---- 左：法人口座 ----
L, LW, T = Inches(0.2), Inches(6.0), Inches(1.1)

tb(slide, "法人口座の作り方",
   L, T, LW, Inches(0.4), font_size=12, bold=True, color=COLOR_BLACK)

flow_items = [
    ("入金用口座",     COLOR_BLUE,   "売上・入金の受け取り先。法人口座のメイン。"),
    ("支払用口座",     COLOR_ORANGE, "経費・人件費などの支払い専用。毎月入金口座から移動。"),
    ("納税用口座",     COLOR_RED,    "法人税・消費税の積立専用。決して取り崩さない。"),
    ("設備投資用口座", COLOR_GREEN,  "大型購入・将来の投資のための積立口座。"),
]
fh = Inches(0.78)
for i, (title, color, desc) in enumerate(flow_items):
    fy = T + Inches(0.5) + i * (fh + Inches(0.06))
    add_rounded_rect(slide, L, fy, LW, fh, fill_color=COLOR_WHITE, radius=0.04)
    add_rect(slide, L, fy, Inches(0.2), fh, fill_color=color)
    tb(slide, title, L + Inches(0.3), fy + Inches(0.08), LW - Inches(0.35), Inches(0.32),
       font_size=10, bold=True, color=COLOR_BLACK)
    tb(slide, desc,  L + Inches(0.3), fy + Inches(0.42), LW - Inches(0.35), Inches(0.3),
       font_size=8, color=COLOR_DARK_GRAY)

bank_top = T + Inches(0.5) + 4 * (fh + Inches(0.06)) + Inches(0.12)
bank_cats = [
    ("都市銀行",   COLOR_BLUE,  "三井住友銀行 Trunk",         "他行宛145円・月額無料\n融資・補助金・対外信用◎"),
    ("ネット銀行", COLOR_GREEN, "GMOあおぞらネット銀行",       "他行宛130円（業界最安）\nPay-easy対応・翌営業日開設"),
    ("地方銀行",   COLOR_WARM,  "地元の信用金庫（任意）",      "融資関係が目的の場合のみ\n日常口座には不向き"),
]
bw = (LW - Inches(0.16)) / 3
bh = Inches(1.32)
for i, (cat, color, name, note) in enumerate(bank_cats):
    bx = L + i * (bw + Inches(0.08))
    add_rounded_rect(slide, bx, bank_top, bw, bh, fill_color=COLOR_WHITE, radius=0.04)
    add_rect(slide, bx, bank_top, bw, Inches(0.32), fill_color=color)
    tb(slide, cat, bx, bank_top, bw, Inches(0.32),
       font_size=9, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    tb(slide, name, bx + Cm(0.2), bank_top + Inches(0.38), bw - Cm(0.4), Inches(0.32),
       font_size=8.5, bold=True, color=COLOR_BLACK)
    tb(slide, note, bx + Cm(0.2), bank_top + Inches(0.74), bw - Cm(0.4), Inches(0.5),
       font_size=7.5, color=COLOR_DARK_GRAY)

# ---- 右：法人クレジットカード ----
R, RW = Inches(6.5), Inches(6.6)

tb(slide, "法人クレジットカード",
   R, T, RW, Inches(0.4), font_size=12, bold=True, color=COLOR_BLACK)

cards = [
    {
        "name": "MFビジネスカード",
        "tag": "当事務所 第1推奨",
        "color": COLOR_GREEN,
        "merit": "MFクラウド会計とAPI完全連携。利用→明細自動取得→仕訳自動生成が途切れない。年会費実質無料・個人事業主も申込可。",
        "fit": "MFクラウドを使う方 全員に推奨（当事務所推奨）",
    },
    {
        "name": "UPSIDER",
        "tag": "スタートアップ向け・高還元",
        "color": COLOR_ORANGE,
        "merit": "還元率1〜1.5%・決算書不要で審査最短10分。設立直後の法人でも通過率が高く、限度額は最大10億円。",
        "fit": "設立間もない法人・審査が不安な方",
    },
    {
        "name": "バクラク",
        "tag": "経理DX・中規模企業向け",
        "color": COLOR_PURPLE,
        "merit": "キャッシュバック1%（半年間1.5%）。請求書・経費精算・支払管理を一元化。電子帳簿保存法対応済み。",
        "fit": "経費精算をシステム化・組織化したい方",
    },
    {
        "name": "AMEX（参考）",
        "tag": "海外出張・ステータス用途のみ",
        "color": COLOR_BLUE,
        "merit": "ブランド力・旅行保険・ラウンジは充実。ただしMFクラウドとのAPI連携が非対応・年会費3.6万円〜のため主軸カードには非推奨。",
        "fit": "海外出張が多い方のサブカードとして",
    },
]

card_h = Inches(1.26)
card_gap = Inches(0.06)
for i, card in enumerate(cards):
    cy = T + Inches(0.5) + i * (card_h + card_gap)
    add_rounded_rect(slide, R, cy, RW, card_h, fill_color=COLOR_WHITE, radius=0.03)
    add_rect(slide, R, cy, Inches(0.2), card_h, fill_color=card["color"])
    tag_w = Inches(2.3)
    add_rect(slide, R + Inches(0.3), cy + Inches(0.09), tag_w, Cm(0.26), fill_color=card["color"])
    tb(slide, card["tag"],
       R + Inches(0.3), cy + Inches(0.07), tag_w, Cm(0.28),
       font_size=7, bold=True, color=COLOR_WHITE)
    tb(slide, card["name"],
       R + Inches(0.3), cy + Inches(0.40), RW - Inches(0.5), Inches(0.34),
       font_size=13, bold=True, color=COLOR_BLACK)
    tb(slide, card["merit"],
       R + Inches(0.3), cy + Inches(0.76), RW - Inches(0.5), Inches(0.34),
       font_size=8, color=COLOR_DARK_GRAY)
    tb(slide, f"→ {card['fit']}",
       R + Inches(0.3), cy + Inches(1.10), RW - Inches(0.5), Inches(0.18),
       font_size=7.5, bold=True, color=card["color"])

footer(slide, "※ 口座開設・クレジットカード発行の代行も承ります。口座1行：30,000円 / 追加口座またはカード発行：各10,000円（すべて税抜）")


# ============================================================
# Slides 17-18: 08 SEVENRICHとのやり取り＆SOSリスト
# ============================================================
section_door(prs, "08", "SEVENRICHとのやり取り＆SOSリスト",
             "連絡方法・資料提出のルール・動く前に必ず相談すること",
             color=COLOR_BLUE)

slide = prs.slides.add_slide(BL)
add_bg(slide, BG_SLIDE)
header(slide, "08", "SEVENRICHとのやり取り＆SOSリスト")

# 左：コミュニケーションルール
L, LW, T = Inches(0.2), Inches(6.0), Inches(1.1)

tb(slide, "コミュニケーションルール",
   L, T, LW, Inches(0.4), font_size=12, bold=True, color=COLOR_BLACK)

comm_rules = [
    ("連絡手段", COLOR_BLUE,
     "専用の Slack チャンネルからいつでもお気軽にご連絡ください。\nご入社後に招待リンクをお送りいたします。"),
    ("資料の連携", COLOR_ORANGE,
     "前月末締めの通帳データ・領収書・請求書を、翌月5日までに\nMFクラウドまたは Slack にてご共有ください（例：締め日の目安）。"),
    ("返信の目安", COLOR_GREEN,
     "通常のメッセージには2営業日以内にご返信いたします。\n税務上の判断が必要な場合は、お時間をいただく場合がございます。"),
    ("月次面談", COLOR_RED,
     "定期的に試算表をもとに確認ミーティングを実施いたします。\n詳細は契約内容をご確認ください。"),
]

for i, (title, color, desc) in enumerate(comm_rules):
    ry = T + Inches(0.5) + i * Inches(1.3)
    add_rounded_rect(slide, L, ry, LW, Inches(1.18), fill_color=COLOR_WHITE, radius=0.03)
    add_rect(slide, L, ry, Inches(0.18), Inches(1.18), fill_color=color)
    tb(slide, title,
       L + Inches(0.28), ry + Inches(0.1), LW - Inches(0.4), Inches(0.35),
       font_size=10, bold=True, color=COLOR_BLACK)
    tb(slide, desc,
       L + Inches(0.28), ry + Inches(0.5), LW - Inches(0.4), Inches(0.6),
       font_size=8.5, color=COLOR_DARK_GRAY)

# 右：SOSリスト
R, RW = Inches(6.5), Inches(6.6)

tb(slide, "ご実施前に必ずひとことご相談ください！ SOSリスト",
   R, T, RW, Inches(0.4), font_size=12, bold=True, color=COLOR_BLACK)

add_rect(slide, R, T + Inches(0.5), RW, Inches(0.45), fill_color=RGBColor(0xFF, 0xEE, 0xEE))
tb(slide, "事後のご報告では税務上のメリットを受けられないケースがございます。必ず「ご実施前に」チャットでひとことご連絡ください。",
   R + Inches(0.1), T + Inches(0.52), RW - Inches(0.2), Inches(0.38),
   font_size=8.5, bold=True, color=COLOR_RED)

sos_items = [
    ("役員報酬を変更したい",           "変更タイミングが損金算入の可否を決めます"),
    ("従業員を雇うことになった",        "社会保険・雇用保険の手続きが必要となります"),
    ("オフィス・登記住所を移転する",    "税務署・自治体・法務局への届出が必要となります"),
    ("100万円以上の大きな買い物をした", "資産計上か経費かの判断が必要となります"),
    ("銀行から借入・出資を受けたい",    "返済計画と税務の両面からアドバイスいたします"),
    ("会社名・事業内容を変更したい",    "定款変更・登記・税務届出が伴います"),
    ("役員を追加・変更したい",          "株主総会議事録・社会保険の手続きが必要となります"),
    ("廃業・会社清算を考えている",      "タイミングにより税負担が大きく変わります"),
    ("大きな投資や保険の加入を検討中",  "節税効果のシミュレーションを事前にご提案いたします"),
]

for i, (t, d) in enumerate(sos_items):
    row = i % 5
    col = i // 5
    sx = R + col * Inches(3.3)
    sy = T + Inches(1.1) + row * Inches(1.0)
    add_rounded_rect(slide, sx, sy, Inches(3.1), Inches(0.88), fill_color=COLOR_WHITE, radius=0.04)
    add_rect(slide, sx, sy, Inches(0.18), Inches(0.88), fill_color=COLOR_RED)
    tb(slide, t,
       sx + Inches(0.28), sy + Inches(0.08), Inches(2.7), Inches(0.35),
       font_size=9.5, bold=True, color=COLOR_BLACK)
    tb(slide, d,
       sx + Inches(0.28), sy + Inches(0.48), Inches(2.7), Inches(0.35),
       font_size=7.5, color=COLOR_DARK_GRAY)

footer(slide, "「こんなことでも相談していいのかな？」と思われたことでも、どうぞお気軽にご連絡ください。早めのご相談ほど、より多くの選択肢をご提案できます。")


# ============================================================
# Slides 19-20: 09 FAQ
# ============================================================
section_door(prs, "09", "よくある質問（FAQ）",
             "新規顧客からよく頂く質問と回答をまとめました")

slide = prs.slides.add_slide(BL)
add_bg(slide, BG_SLIDE)
header(slide, "09", "よくある質問（FAQ）")

faqs = [
    {
        "q": "会計ソフトは何を使えばいいですか？",
        "a": "当事務所ではMFクラウド会計を推奨しております。銀行・クレカと自動連携でき、月次仕訳を大幅に効率化できます。初期設定は当事務所が一緒に行いますのでご安心ください。",
        "color": COLOR_RED,
    },
    {
        "q": "領収書や請求書はどう管理すればいいですか？",
        "a": "電子取引のデータは原則データでの保存が必要です（電子帳簿保存法）。紙の領収書はMFクラウドでの写真保存が便利です。まずは「捨てない」「専用の箱に入れる」ことから始めてみてください。",
        "color": COLOR_ORANGE,
    },
    {
        "q": "消費税の免税はいつまで続きますか？",
        "a": "原則、設立後2期は免税となります（資本金1,000万円未満・特定期間売上1,000万円以下の場合）。ただしインボイス登録をされると、免税期間中でも課税事業者となりますのでご注意ください。",
        "color": COLOR_ORANGE,
    },
    {
        "q": "役員報酬はいつでも変更できますか？",
        "a": "定期同額給与として損金算入するには、事業年度開始から3ヶ月以内に変更する必要があります。期間外の変更分は損金不算入となります。変更をご検討の際は、必ず事前にご相談ください。",
        "color": COLOR_RED,
    },
    {
        "q": "赤字でも税金はかかりますか？",
        "a": "法人税・事業税は所得ゼロであれば課税されませんが、法人住民税の均等割（年約7万円〜）は赤字でも発生いたします。また消費税は売上・仕入の状況により別途課税される場合がございます。",
        "color": COLOR_BLUE,
    },
    {
        "q": "SEVENRICHには何をお任せできますか？",
        "a": "月次仕訳チェック・試算表作成・税務相談・決算申告・給与計算サポートなどをご提供しております。経営判断に関わる数字のご相談も、いつでもお気軽にお声がけください。",
        "color": COLOR_GREEN,
    },
]

faq_lefts = [Inches(0.2), Inches(6.6)]
faq_top = Inches(1.1)
col_w_faq = Inches(6.1)
col_h_faq = Inches(1.72)

for i, faq in enumerate(faqs):
    col = i % 2
    row = i // 2
    fl = faq_lefts[col]
    ft = faq_top + row * (col_h_faq + Inches(0.06))

    add_rounded_rect(slide, fl, ft, col_w_faq, col_h_faq, fill_color=COLOR_WHITE, radius=0.03)
    add_rect(slide, fl, ft, Inches(0.18), col_h_faq, fill_color=faq["color"])

    add_rect(slide, fl + Inches(0.28), ft + Inches(0.1),
             Cm(0.55), Cm(0.55), fill_color=faq["color"])
    tb(slide, "Q", fl + Inches(0.28), ft + Inches(0.09), Cm(0.55), Cm(0.56),
       font_size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    tb(slide, faq["q"],
       fl + Inches(0.85), ft + Inches(0.1), col_w_faq - Inches(1.0), Inches(0.45),
       font_size=10, bold=True, color=COLOR_BLACK)

    add_rect(slide, fl + Inches(0.28), ft + Inches(0.73),
             Cm(0.55), Cm(0.55), fill_color=COLOR_DARK_GRAY)
    tb(slide, "A", fl + Inches(0.28), ft + Inches(0.72), Cm(0.55), Cm(0.56),
       font_size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    tb(slide, faq["a"],
       fl + Inches(0.85), ft + Inches(0.73), col_w_faq - Inches(1.0), Inches(0.85),
       font_size=8.5, color=COLOR_DARK_GRAY)

footer(slide, "ご不明な点は何でもお気軽にご相談ください。  SEVENRICH 会計事務所  ｜  専用 Slack チャンネルよりご連絡ください")


# ============================================================
# Closing Slide
# ============================================================
slide = prs.slides.add_slide(BL)
add_bg(slide, COLOR_BLACK)

add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill_color=COLOR_RED)

tb(slide, "ともに、成長してまいりましょう。",
   Inches(1.2), Inches(2.0), Inches(11.0), Inches(1.0),
   font_size=32, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
tb(slide, "ご不明な点・ご相談はいつでもお気軽にお声がけください。\n些細なことでも、早めのご相談ほどより多くのご提案ができます。",
   Inches(1.5), Inches(3.1), Inches(10.3), Inches(0.8),
   font_size=12, color=COLOR_MID_GRAY, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(3.3), Inches(4.3), Inches(6.7), Inches(1.8),
         fill_color=RGBColor(0x28, 0x28, 0x28))
add_rect(slide, Inches(3.3), Inches(4.3), Inches(6.7), Inches(0.06),
         fill_color=COLOR_RED)
tb(slide, "SEVENRICH 会計事務所",
   Inches(3.3), Inches(4.4), Inches(6.7), Inches(0.45),
   font_size=15, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
tb(slide, "専用 Slack チャンネルよりお気軽にご連絡ください",
   Inches(3.3), Inches(4.95), Inches(6.7), Inches(0.4),
   font_size=11, color=COLOR_MID_GRAY, align=PP_ALIGN.CENTER)
tb(slide, "2026年版",
   Inches(0.3), Inches(7.0), Inches(3.0), Inches(0.35),
   font_size=9, color=RGBColor(0x55, 0x55, 0x55))


# ============================================================
# 保存
# ============================================================
output_path = r"C:\Users\HirotsuguTomari\dev\新規顧客ガイドブック\新規顧客ガイドブック_骨子v12.pptx"
prs.save(output_path)
print(f"保存完了: {output_path}")
